#!/usr/bin/env python3
"""
文档解析模块

将支持的 Office 文档和文本型 PDF 文件转换为 JSON Markdown 格式。
支持的文件类型：
- .docx (Word 文档)
- .pdf  (PDF 文件，需要包含可提取的文本)
- .pptx (PowerPoint 演示文稿)

输出格式：JSON 对象，包含：
- markdown: 转换后的 Markdown 文本
- pages: 页面信息列表（仅 PDF）
- metadata: 文件元数据（文件类型等）
"""

import argparse
import json
import sys
from pathlib import Path


def fail(code: str, message: str) -> None:
    """
    输出错误信息并退出程序

    Args:
        code: 错误代码，用于程序化处理
        message: 人类可读的错误描述

    输出格式：JSON 对象 {"error_code": "...", "error_message": "..."}
    """
    print(json.dumps({"error_code": code, "error_message": message}, ensure_ascii=False))
    raise SystemExit(2)


def parse_docx(path: Path) -> str:
    """
    解析 Word 文档 (.docx) 并转换为 Markdown 格式

    Args:
        path: Word 文档的文件路径

    Returns:
        str: 转换后的 Markdown 文本

    处理逻辑：
    1. 提取所有段落文本，识别标题样式并转换为 Markdown 标题
    2. 提取所有表格并转换为 Markdown 表格格式
    3. 段落之间用空行分隔
    """
    from docx import Document

    document = Document(path)
    blocks: list[str] = []
    # 遍历所有段落，提取文本内容
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue  # 跳过空段落
        # 检查段落样式，识别标题（Heading 1, Heading 2 等）
        style = (paragraph.style.name or "").lower()
        if style.startswith("heading"):
            # 从样式名称中提取标题级别（数字）
            level = next((char for char in style if char.isdigit()), "1")
            blocks.append(f"{'#' * int(level)} {text}")  # 转换为 Markdown 标题语法
        else:
            blocks.append(text)
    # 处理文档中的所有表格
    for table in document.tables:
        # 将表格转换为 Markdown 表格格式
        # 每行的单元格内容用 " | " 分隔，换行符替换为空格
        rows = [[cell.text.strip().replace("\n", " ") for cell in row.cells] for row in table.rows]
        # 过滤掉完全为空的行
        rows = [row for row in rows if any(row)]
        if rows:
            # 第一行作为表头
            blocks.append("| " + " | ".join(rows[0]) + " |")
            # 添加 Markdown 表格的分隔行
            blocks.append("| " + " | ".join("---" for _ in rows[0]) + " |")
            # 添加剩余的数据行
            blocks.extend("| " + " | ".join(row) + " |" for row in rows[1:])
    return "\n\n".join(blocks).strip()


def parse_pdf(path: Path) -> tuple[str, list[dict[str, object]]]:
    """
    解析 PDF 文件并转换为 Markdown 格式

    Args:
        path: PDF 文件的文件路径

    Returns:
        tuple: (Markdown 文本, 页面信息列表)
            - Markdown 文本：每页内容用 page 注释分隔
            - 页面信息列表：包含每页的页码和文本内容

    处理逻辑：
    1. 使用 pdfplumber 提取每页的文本内容
    2. 为每页添加 HTML 注释标记（<!-- page:N -->）
    3. 检查提取的文本是否足够（少于20字符视为无有效文本）
    4. 如果缺少可提取文本，提示需要 OCR 处理
    """
    import pdfplumber

    pages: list[dict[str, object]] = []
    blocks: list[str] = []
    # 打开 PDF 文件并遍历每一页
    with pdfplumber.open(path) as pdf:
        for number, page in enumerate(pdf.pages, start=1):
            # 提取当前页的文本内容
            text = (page.extract_text() or "").strip()
            # 记录页面信息（页码和文本）
            pages.append({"number": number, "text": text})
            if text:
                # 添加页面标记注释和文本内容
                blocks.extend((f"<!-- page:{number} -->", text))
    markdown = "\n\n".join(blocks).strip()
    # 检查提取的文本是否足够（去除空白后少于20字符视为无有效文本）
    if len("".join(markdown.split())) < 20:
        fail("PARSER_OCR_REQUIRED", "PDF 缺少可提取文本，当前未启用 OCR")
    return markdown, pages


def parse_pptx(path: Path) -> str:
    """
    解析 PowerPoint 演示文稿 (.pptx) 并转换为 Markdown 格式

    Args:
        path: PowerPoint 文件的文件路径

    Returns:
        str: 转换后的 Markdown 文本

    处理逻辑：
    1. 遍历每张幻灯片，添加页码标题
    2. 提取幻灯片中的文本框内容
    3. 提取幻灯片中的表格并转换为 Markdown 表格
    4. 每张幻灯片之间用空行分隔
    """
    from pptx import Presentation

    blocks: list[str] = []
    # 遍历所有幻灯片
    for number, slide in enumerate(Presentation(path).slides, start=1):
        # 为每张幻灯片添加标题
        blocks.append(f"# 第 {number} 页")
        # 遍历幻灯片中的所有形状（文本框、表格、图片等）
        for shape in slide.shapes:
            # 处理文本框：提取文本内容
            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                if text:
                    blocks.append(text)
            # 处理表格：转换为 Markdown 表格格式
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                    blocks.append("| " + " | ".join(cells) + " |")
    return "\n\n".join(blocks).strip()


def main() -> None:
    """
    主函数：解析命令行参数并执行文档解析

    命令行参数：
        --input: 必需，输入文件的路径

    处理流程：
    1. 解析命令行参数
    2. 验证输入文件是否存在
    3. 根据文件扩展名选择对应的解析器
    4. 输出 JSON 格式的解析结果
    """
    # 设置命令行参数解析器
    parser = argparse.ArgumentParser()
    parser.add_argument("--input", required=True)
    args = parser.parse_args()
    path = Path(args.input)
    # 验证输入文件是否存在
    if not path.is_file():
        fail("PARSER_FILE_NOT_FOUND", "输入文件不存在")
    # 获取文件扩展名并转换为小写
    suffix = path.suffix.lower()
    try:
        # 根据文件类型调用对应的解析器
        if suffix == ".docx":
            markdown, pages = parse_docx(path), []  # Word 文档不需要页面信息
        elif suffix == ".pdf":
            markdown, pages = parse_pdf(path)  # PDF 需要返回页面信息
        elif suffix == ".pptx":
            markdown, pages = parse_pptx(path), []  # PPT 不需要页面信息
        else:
            fail("PARSER_UNSUPPORTED_TYPE", f"不支持的文档类型: {suffix}")
    except SystemExit:
        raise  # 重新抛出 SystemExit（来自 fail() 函数）
    except Exception as exc:
        # 捕获其他所有异常，记录错误并输出错误信息
        print(f"parser failed: {exc}", file=sys.stderr)
        fail("PARSER_FAILED", "文档解析失败")
    # 输出 JSON 格式的解析结果
    print(json.dumps({"markdown": markdown, "pages": pages, "metadata": {"file_type": suffix}}, ensure_ascii=False))


if __name__ == "__main__":
    main()  # 脚本入口点
